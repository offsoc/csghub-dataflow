import os
import shutil
import time
import traceback
from data_celery.main import celery_app
from sqlalchemy.orm import Session
from loguru import logger
from pathlib import Path

from data_celery.utils import get_format_folder_path, ensure_directory_exists, ensure_directory_exists_remove
from data_celery.db import FormatifyManager
from data_engine.exporter.load import load_exporter
from data_server.formatify.FormatifyModels import DataFormatTask, DataFormatTaskStatusEnum, DataFormatTypeEnum, \
    getFormatTypeName
from data_server.database.session import get_sync_session
from data_engine.ingester.load import load_ingester
import pandas as pd
import mammoth
from markdownify import markdownify as md
from pptx import Presentation
from data_celery.mongo_tools.tools import insert_formatity_task_log_info, insert_formatity_task_log_error
from pycsghub.upload_large_folder.main import upload_large_folder_internal
from pycsghub.cmd.repo_types import RepoType
from pycsghub.utils import (build_csg_headers,
                            model_id_to_group_owner_name,
                            get_endpoint,
                            REPO_TYPE_DATASET)
from data_engine.utils.env import GetHubEndpoint
@celery_app.task
def format_task(task_id: int, user_name: str, user_token: str):
    tmp_path: str = None
    db_session: Session = None
    format_task: DataFormatTask = None
    try:
        db_session: Session = get_sync_session()
        format_task: DataFormatTask = FormatifyManager.get_formatify_task(db_session, task_id)
        tmp_path = get_format_folder_path(format_task.task_uid)
        insert_formatity_task_log_info(format_task.task_uid, f"创建临时目录：{tmp_path}")
        ensure_directory_exists(tmp_path)
        # 下载源目录
        insert_formatity_task_log_info(format_task.task_uid, f"开始下载源目录...")
        ingesterCSGHUB = load_ingester(
            dataset_path=tmp_path,
            repo_id=format_task.from_csg_hub_repo_id,
            branch=format_task.from_csg_hub_dataset_branch,
            user_name=user_name,
            user_token=user_token,
        )
        ingester_result = ingesterCSGHUB.ingest()
        insert_formatity_task_log_info(format_task.task_uid, f"下载源目录完成...目录地址：{ingester_result}")
        work_dir = Path(tmp_path).joinpath('work')
        insert_formatity_task_log_info(format_task.task_uid, f"开始转换文件...")
        # 转换
        format_task_func(
            tmp_path=ingester_result,
            from_type=format_task.from_data_type,
            to_type=format_task.to_data_type,
            task_uid=format_task.task_uid,
        )
        insert_formatity_task_log_info(format_task.task_uid, f"转换文件完成...")
        insert_formatity_task_log_info(format_task.task_uid, f"开始上传目标目录...")
        # 上传目标目录
        exporter = load_exporter(
            export_path=ingester_result,
            repo_id=format_task.to_csg_hub_repo_id,
            branch=format_task.to_csg_hub_dataset_default_branch,
            user_name=user_name,
            user_token=user_token,
            path_is_dir=True,
            work_dir=str(work_dir)
        )
        exporter.export_large_folder()
        insert_formatity_task_log_info(format_task.task_uid, '上传完成...')
        format_task.task_status = DataFormatTaskStatusEnum.COMPLETED.value
        db_session.commit()
        pass
    except Exception as e:
        traceback.print_exc()
        format_task.task_status = DataFormatTaskStatusEnum.ERROR.value
        db_session.commit()
        insert_formatity_task_log_error(format_task.task_uid, f"转换任务失败: {str(e)}")
    finally:
        pass
        # 删除tmp_path
        if tmp_path:
            shutil.rmtree(tmp_path)
            insert_formatity_task_log_info(format_task.task_uid, f"删除临时目录：{tmp_path}")


def format_task_func(
        tmp_path: str,
        from_type: int,
        to_type: int,
        task_uid: str
):
    insert_formatity_task_log_info(task_uid,
                                   f"转换目录：{tmp_path}，源文件类型：{getFormatTypeName(from_type)}，目标文件类型：{getFormatTypeName(to_type)}")
    match from_type:
        case DataFormatTypeEnum.Excel.value:
            match to_type:
                case DataFormatTypeEnum.Csv.value:
                    traverse_files(tmp_path, convert_excel_to_csv, task_uid)
                case DataFormatTypeEnum.Json.value:
                    traverse_files(tmp_path, convert_excel_to_json, task_uid)
                case DataFormatTypeEnum.Parquet.value:
                    traverse_files(tmp_path, convert_excel_to_parquet, task_uid)
        case DataFormatTypeEnum.Word.value:
            match to_type:
                case DataFormatTypeEnum.Markdown.value:
                    traverse_files(tmp_path, convert_word_to_markdown, task_uid)
        case DataFormatTypeEnum.PPT.value:
            match to_type:
                case DataFormatTypeEnum.Markdown.value:
                    traverse_files(tmp_path, convert_ppt_to_markdown, task_uid)


def traverse_files(file_path: str, func, task_uid):
    for root, dirs, files in os.walk(file_path):
        for file in files:
            file_path = os.path.join(root, file)
            func(file_path, task_uid)
        for dir in dirs:
            dir_path = os.path.join(root, dir)
            traverse_files(dir_path, func, task_uid)


def convert_excel_to_csv(file_path: str, task_uid):
    if file_path.lower().endswith(('.xlsx', '.xls')):
        insert_formatity_task_log_info(task_uid, f'源文件地址：{file_path}')
        try:
            df = pd.read_excel(file_path)
            new_file = os.path.splitext(file_path)[0] + '.csv'
            df.to_csv(new_file, index=False)
            insert_formatity_task_log_info(task_uid, f'转换文件 {new_file} 成功')
            os.remove(file_path)
            return True
        except Exception as e:
            print(f"转换文件 {file_path} 时出错: {e}")
            insert_formatity_task_log_error(task_uid, f"转换文件 {file_path} 时出错: {e}")
            return False
    else:
        return True


def convert_excel_to_json(file_path: str, task_uid):
    if file_path.lower().endswith(('.xlsx', '.xls')):
        insert_formatity_task_log_info(task_uid, f'源文件地址：{file_path}')
        try:
            df = pd.read_excel(file_path)
            new_file = os.path.splitext(file_path)[0] + '.json'
            df.to_json(new_file, orient='records', force_ascii=False)
            insert_formatity_task_log_info(task_uid, f'转换文件 {new_file} 成功')
            os.remove(file_path)
            return True
        except Exception as e:
            print(f"转换文件 {file_path} 时出错: {e}")
            insert_formatity_task_log_error(task_uid, f"转换文件 {file_path} 时出错: {e}")

            return False
    else:
        # 非 Excel 文件，跳过
        return True


def convert_excel_to_parquet(file_path: str, task_uid):
    if file_path.lower().endswith(('.xlsx', '.xls')):
        insert_formatity_task_log_info(task_uid, f'源文件地址：{file_path}')
        try:
            df = pd.read_excel(file_path)
            new_file = os.path.splitext(file_path)[0] + '.parquet'
            df.to_parquet(new_file + '.parquet', index=False)
            insert_formatity_task_log_info(task_uid, f'转换文件 {new_file} 成功')
            os.remove(file_path)
            return True
        except Exception as e:
            print(f"转换文件 {file_path} 时出错: {e}")
            insert_formatity_task_log_error(task_uid, f"转换文件 {file_path} 时出错: {e}")
            return False
    else:
        return True


def convert_word_to_markdown(file_path: str, task_uid):
    if file_path.lower().endswith(('.docx', '.doc')):
        insert_formatity_task_log_info(task_uid, f'源文件地址：{file_path}')
        try:
            with open(file_path, "rb") as docx_file:
                result = mammoth.convert_to_html(docx_file)
                html_content = result.value
            markdown_content = md(html_content)
            markdown_file_path = os.path.splitext(file_path)[0] + '.md'
            with open(markdown_file_path, 'w', encoding='utf-8') as md_file:
                md_file.write(markdown_content)
            insert_formatity_task_log_info(task_uid, f'转换文件 {markdown_file_path} 成功')
            os.remove(file_path)
            return True
        except Exception as e:
            print(f"转换文件 {file_path} 时出错: {e}")
            insert_formatity_task_log_error(task_uid, f"转换文件 {file_path} 时出错: {e}")
            return False
    else:
        # 非 Word 文件，跳过
        return True


def convert_ppt_to_markdown(file_path: str, task_uid):
    if file_path.lower().endswith(('.pptx', '.ppt')):
        insert_formatity_task_log_info(task_uid, f'源文件地址：{file_path}')
        try:
            prs = Presentation(file_path)
            markdown_content = ""
            for i, slide in enumerate(prs.slides):
                markdown_content += f"# 幻灯片 {i + 1}\n\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content = shape.text.strip()
                        if text_content:
                            if len(text_content) < 50 and '\n' not in text_content:
                                markdown_content += f"## {text_content}\n\n"
                            else:
                                markdown_content += f"{text_content}\n\n"
            markdown_file_path = os.path.splitext(file_path)[0] + '.md'
            with open(markdown_file_path, 'w', encoding='utf-8') as md_file:
                md_file.write(markdown_content)
            insert_formatity_task_log_info(task_uid, f'转换文件 {markdown_file_path} 成功')
            os.remove(file_path)
            return True
        except Exception as e:
            print(f"转换文件 {file_path} 时出错: {e}")
            insert_formatity_task_log_error(task_uid, f"转换文件 {file_path} 时出错: {e}")
            return False
    else:
        # 非 PPT 文件，跳过
        return True
